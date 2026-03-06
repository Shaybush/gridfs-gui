"""Document conversion service for file preview.

Converts office documents to HTML (via python-docx, python-pptx, openpyxl),
CSV to styled HTML tables, and Markdown to syntax-highlighted HTML.
If LibreOffice is available, it is used for higher fidelity PDF conversion.
All conversions produce self-contained output suitable for embedding in an
iframe or returning directly to the client.
"""

import asyncio
import csv
import html
import io
import logging
import math
import os
import shutil
import tempfile

import markdown as md
from pygments.formatters import HtmlFormatter

from app.config import get_settings

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# MIME-type / extension mappings
# ---------------------------------------------------------------------------

OFFICE_MIME_TYPES: dict[str, str] = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/msword": ".doc",
    "application/vnd.ms-powerpoint": ".ppt",
    "application/vnd.ms-excel": ".xls",
    "application/vnd.oasis.opendocument.text": ".odt",
    "application/vnd.oasis.opendocument.spreadsheet": ".ods",
    "application/vnd.oasis.opendocument.presentation": ".odp",
}

OFFICE_EXTENSIONS: set[str] = {
    ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".odt", ".ods", ".odp",
}

# Extensions that can be handled by pure Python (no LibreOffice needed)
_PYTHON_CONVERTIBLE: set[str] = {".docx", ".pptx", ".xlsx"}

CSV_MIME_TYPES: set[str] = {"text/csv"}
MARKDOWN_MIME_TYPES: set[str] = {"text/markdown"}

# ---------------------------------------------------------------------------
# Shared HTML styles
# ---------------------------------------------------------------------------

_BASE_STYLES = """
* { box-sizing: border-box; margin: 0; padding: 0; }
body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       padding: 20px; background: #1e1e1e; color: #d4d4d4; line-height: 1.6; }
"""

_LIBREOFFICE_AVAILABLE: bool | None = None


def _has_libreoffice() -> bool:
    global _LIBREOFFICE_AVAILABLE
    if _LIBREOFFICE_AVAILABLE is None:
        _LIBREOFFICE_AVAILABLE = shutil.which("soffice") is not None
    return _LIBREOFFICE_AVAILABLE


class DocumentConverter:
    """Stateless converter with methods for office, CSV, and Markdown files."""

    # ------------------------------------------------------------------
    # Office -> PDF via LibreOffice (when available)
    # ------------------------------------------------------------------

    @staticmethod
    async def convert_office_to_pdf(file_bytes: bytes, source_extension: str) -> bytes:
        """Convert an office document to PDF using LibreOffice headless mode."""
        settings = get_settings()
        timeout = settings.PREVIEW_CONVERSION_TIMEOUT_SECONDS

        tmp_dir = tempfile.mkdtemp(prefix="gridfs_preview_")
        src_path = os.path.join(tmp_dir, f"source{source_extension}")

        try:
            with open(src_path, "wb") as f:
                f.write(file_bytes)

            env = os.environ.copy()
            env["HOME"] = "/tmp/libreoffice-profile"

            process = await asyncio.create_subprocess_exec(
                "soffice",
                "--headless",
                "--norestore",
                f"-env:UserInstallation=file:///tmp/libreoffice-profile",
                "--convert-to", "pdf",
                "--outdir", tmp_dir,
                src_path,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                env=env,
            )

            try:
                stdout, stderr = await asyncio.wait_for(
                    process.communicate(), timeout=timeout,
                )
            except asyncio.TimeoutError:
                process.kill()
                await process.wait()
                raise RuntimeError(
                    f"LibreOffice conversion timed out after {timeout}s"
                )

            if process.returncode != 0:
                err_msg = stderr.decode("utf-8", errors="replace").strip()
                raise RuntimeError(
                    f"LibreOffice exited with code {process.returncode}: {err_msg}"
                )

            pdf_path = os.path.join(tmp_dir, "source.pdf")
            if not os.path.isfile(pdf_path):
                raise FileNotFoundError(
                    f"Expected PDF output not found at {pdf_path}"
                )

            with open(pdf_path, "rb") as f:
                return f.read()

        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    # ------------------------------------------------------------------
    # Office -> HTML via pure Python libraries
    # ------------------------------------------------------------------

    @staticmethod
    def convert_office_to_html(file_bytes: bytes, source_extension: str) -> str:
        """Convert an office document to HTML using pure Python libraries.

        Supports .docx, .pptx, and .xlsx. Falls back to a basic message
        for unsupported formats (.doc, .ppt, .xls, .odt, etc.).
        """
        ext = source_extension.lower()
        if ext == ".docx":
            return DocumentConverter._convert_docx_to_html(file_bytes)
        elif ext == ".pptx":
            return DocumentConverter._convert_pptx_to_html(file_bytes)
        elif ext == ".xlsx":
            return DocumentConverter._convert_xlsx_to_html(file_bytes)
        else:
            raise RuntimeError(
                f"Preview for '{ext}' files requires LibreOffice which is not installed. "
                f"Supported formats without LibreOffice: .docx, .pptx, .xlsx"
            )

    @staticmethod
    def _convert_docx_to_html(file_bytes: bytes) -> str:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        body_parts: list[str] = []

        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                # Find matching paragraph object
                for para in doc.paragraphs:
                    if para._element is element:
                        body_parts.append(DocumentConverter._render_paragraph(para))
                        break

            elif tag == "tbl":
                for table in doc.tables:
                    if table._element is element:
                        body_parts.append(DocumentConverter._render_table(table))
                        break

        body_html = "\n".join(body_parts)

        return f"""<!DOCTYPE html>
<html lang="en">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<style>
{_BASE_STYLES}
.doc-container {{ max-width: 800px; margin: 0 auto; }}
h1 {{ font-size: 1.8em; color: #e0e0e0; margin: 0.8em 0 0.4em; }}
h2 {{ font-size: 1.5em; color: #e0e0e0; margin: 0.7em 0 0.3em; }}
h3 {{ font-size: 1.3em; color: #e0e0e0; margin: 0.6em 0 0.3em; }}
h4, h5, h6 {{ font-size: 1.1em; color: #e0e0e0; margin: 0.5em 0 0.2em; }}
p {{ margin: 0.4em 0; }}
.bold {{ font-weight: bold; }}
.italic {{ font-style: italic; }}
.underline {{ text-decoration: underline; }}
.strikethrough {{ text-decoration: line-through; }}
table {{ width: 100%; border-collapse: collapse; margin: 1em 0; font-size: 13px; }}
th, td {{ border: 1px solid #444; padding: 8px 10px; text-align: left; }}
th {{ background: #2d2d2d; color: #e0e0e0; }}
tr:nth-child(even) td {{ background: #252525; }}
tr:nth-child(odd) td {{ background: #1e1e1e; }}
ul, ol {{ margin: 0.5em 0 0.5em 1.5em; }}
li {{ margin: 0.2em 0; }}
.indent-1 {{ margin-left: 2em; }}
.indent-2 {{ margin-left: 4em; }}
.indent-3 {{ margin-left: 6em; }}
img {{ max-width: 100%; height: auto; margin: 0.5em 0; }}
</style></head>
<body><div class="doc-container">{body_html}</div></body></html>"""

    @staticmethod
    def _render_paragraph(para) -> str:
        from docx.oxml.ns import qn

        style_name = (para.style.name or "").lower()

        # Determine heading level
        if style_name.startswith("heading"):
            try:
                level = int(style_name.replace("heading", "").strip())
                level = max(1, min(level, 6))
            except ValueError:
                level = None
        else:
            level = None

        # Check for list style
        pPr = para._element.find(qn("w:pPr"))
        is_list = False
        if pPr is not None:
            numId_elem = pPr.find(qn("w:numPr"))
            if numId_elem is not None:
                is_list = True

        # Build inline content with formatting
        parts: list[str] = []
        for run in para.runs:
            text = html.escape(run.text)
            if not text:
                continue
            classes = []
            if run.bold:
                classes.append("bold")
            if run.italic:
                classes.append("italic")
            if run.underline:
                classes.append("underline")
            if run.font.strike:
                classes.append("strikethrough")
            if classes:
                text = f'<span class="{" ".join(classes)}">{text}</span>'
            parts.append(text)

        content = "".join(parts)
        if not content.strip():
            return ""

        if level:
            return f"<h{level}>{content}</h{level}>"

        if is_list:
            return f"<li>{content}</li>"

        # Check indentation
        indent_class = ""
        if para.paragraph_format.left_indent:
            emu = para.paragraph_format.left_indent
            indent_level = min(3, max(1, int(emu / 457200) + 1))
            indent_class = f' class="indent-{indent_level}"'

        return f"<p{indent_class}>{content}</p>"

    @staticmethod
    def _render_table(table) -> str:
        rows_html: list[str] = []
        for i, row in enumerate(table.rows):
            cells = "".join(
                f"<td>{html.escape(cell.text)}</td>" for cell in row.cells
            )
            rows_html.append(f"<tr>{cells}</tr>")
        return f"<table><tbody>{''.join(rows_html)}</tbody></table>"

    @staticmethod
    def _convert_pptx_to_html(file_bytes: bytes) -> str:
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation(io.BytesIO(file_bytes))
        slides_html: list[str] = []

        for idx, slide in enumerate(prs.slides, 1):
            shapes_html: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    paras: list[str] = []
                    for para in shape.text_frame.paragraphs:
                        runs_text: list[str] = []
                        for run in para.runs:
                            text = html.escape(run.text)
                            if not text:
                                continue
                            styles = []
                            if run.font.bold:
                                styles.append("font-weight:bold")
                            if run.font.italic:
                                styles.append("font-style:italic")
                            if run.font.size:
                                pt = run.font.size.pt
                                styles.append(f"font-size:{pt}pt")
                            if styles:
                                text = f'<span style="{";".join(styles)}">{text}</span>'
                            runs_text.append(text)
                        line = "".join(runs_text)
                        if line.strip():
                            paras.append(f"<p>{line}</p>")
                    if paras:
                        shapes_html.append("\n".join(paras))

                elif shape.has_table:
                    table = shape.table
                    rows: list[str] = []
                    for row in table.rows:
                        cells = "".join(
                            f"<td>{html.escape(cell.text)}</td>" for cell in row.cells
                        )
                        rows.append(f"<tr>{cells}</tr>")
                    shapes_html.append(f"<table>{''.join(rows)}</table>")

            content = "\n".join(shapes_html) if shapes_html else "<p><em>Empty slide</em></p>"
            slides_html.append(
                f'<div class="slide"><div class="slide-header">Slide {idx}</div>'
                f'<div class="slide-content">{content}</div></div>'
            )

        all_slides = "\n".join(slides_html)
        total = len(prs.slides)

        return f"""<!DOCTYPE html>
<html lang="en">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<style>
{_BASE_STYLES}
.slide {{ background: #252525; border: 1px solid #444; border-radius: 8px;
         padding: 24px; margin-bottom: 20px; }}
.slide-header {{ font-size: 12px; color: #888; margin-bottom: 12px;
                 padding-bottom: 8px; border-bottom: 1px solid #333; }}
.slide-content p {{ margin: 0.3em 0; }}
.slide-content table {{ width: 100%; border-collapse: collapse; margin: 0.5em 0; }}
.slide-content td {{ border: 1px solid #444; padding: 6px 8px; }}
.info {{ font-size: 13px; color: #888; margin-bottom: 16px; }}
</style></head>
<body>
<div class="info">{total} slide{"s" if total != 1 else ""}</div>
{all_slides}
</body></html>"""

    @staticmethod
    def _convert_xlsx_to_html(file_bytes: bytes) -> str:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sheets_html: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                cells = [html.escape(str(c)) if c is not None else "" for c in row]
                rows.append(cells)

            if not rows:
                sheets_html.append(
                    f'<div class="sheet"><div class="sheet-header">{html.escape(sheet_name)}</div>'
                    f'<p><em>Empty sheet</em></p></div>'
                )
                continue

            # First row as header
            header = rows[0]
            thead = "<tr>" + "".join(f"<th>{c}</th>" for c in header) + "</tr>"
            tbody_parts: list[str] = []
            for i, row in enumerate(rows[1:]):
                cls = "even" if i % 2 == 0 else "odd"
                cells = "".join(f"<td>{c}</td>" for c in row)
                tbody_parts.append(f'<tr class="{cls}">{cells}</tr>')

            table = f"<table><thead>{thead}</thead><tbody>{''.join(tbody_parts)}</tbody></table>"
            sheets_html.append(
                f'<div class="sheet"><div class="sheet-header">{html.escape(sheet_name)}'
                f' <span class="row-count">({len(rows) - 1} rows)</span></div>{table}</div>'
            )

        wb.close()
        all_sheets = "\n".join(sheets_html)

        return f"""<!DOCTYPE html>
<html lang="en">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<style>
{_BASE_STYLES}
.sheet {{ margin-bottom: 24px; }}
.sheet-header {{ font-size: 15px; font-weight: bold; color: #e0e0e0;
                 margin-bottom: 8px; padding-bottom: 6px; border-bottom: 1px solid #444; }}
.row-count {{ font-weight: normal; color: #888; font-size: 12px; }}
table {{ width: 100%; border-collapse: collapse; font-size: 13px; }}
th {{ background: #2d2d2d; color: #e0e0e0; text-align: left;
     padding: 8px 10px; border: 1px solid #444; position: sticky; top: 0; }}
td {{ padding: 6px 10px; border: 1px solid #333; }}
tr.even td {{ background: #1e1e1e; }}
tr.odd td {{ background: #252525; }}
tr:hover td {{ background: #2a2d3e; }}
</style></head>
<body>{all_sheets}</body></html>"""

    # ------------------------------------------------------------------
    # CSV -> HTML table
    # ------------------------------------------------------------------

    @staticmethod
    def convert_csv_to_html(
        file_bytes: bytes,
        page: int = 1,
        rows_per_page: int = 100,
    ) -> tuple[str, int]:
        """Parse CSV bytes and render a paginated, styled HTML table."""
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1")

        reader = csv.reader(io.StringIO(text))
        rows: list[list[str]] = list(reader)

        if not rows:
            return "<p>Empty CSV file.</p>", 1

        header = rows[0]
        data_rows = rows[1:]
        total_pages = max(1, math.ceil(len(data_rows) / rows_per_page))
        page = max(1, min(page, total_pages))

        start = (page - 1) * rows_per_page
        end = start + rows_per_page
        page_rows = data_rows[start:end]

        table_rows = "".join(f"<th>{html.escape(h)}</th>" for h in header)
        thead = f"<tr>{table_rows}</tr>"

        tbody_parts: list[str] = []
        for idx, row in enumerate(page_rows):
            cls = "even" if idx % 2 == 0 else "odd"
            cells = "".join(f"<td>{html.escape(c)}</td>" for c in row)
            tbody_parts.append(f'<tr class="{cls}">{cells}</tr>')
        tbody = "".join(tbody_parts)

        document = f"""<!DOCTYPE html>
<html lang="en">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<style>
{_BASE_STYLES}
.info {{ font-size: 13px; margin-bottom: 8px; color: #999; }}
table {{ width: 100%; border-collapse: collapse; font-size: 13px; }}
th {{ background: #2d2d2d; color: #e0e0e0; text-align: left;
     padding: 8px 10px; border: 1px solid #444; position: sticky; top: 0; }}
td {{ padding: 6px 10px; border: 1px solid #333; }}
tr.even td {{ background: #1e1e1e; }}
tr.odd td  {{ background: #252525; }}
tr:hover td {{ background: #2a2d3e; }}
</style></head>
<body>
<div class="info">Page {page} of {total_pages} &middot; {len(data_rows)} rows total</div>
<table><thead>{thead}</thead><tbody>{tbody}</tbody></table>
</body></html>"""

        return document, total_pages

    # ------------------------------------------------------------------
    # Markdown -> HTML
    # ------------------------------------------------------------------

    @staticmethod
    def convert_markdown_to_html(file_bytes: bytes) -> str:
        """Convert Markdown bytes to self-contained, syntax-highlighted HTML."""
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1")

        extensions = ["tables", "fenced_code", "codehilite", "toc"]
        extension_configs = {
            "codehilite": {
                "css_class": "codehilite",
                "guess_lang": True,
                "noclasses": False,
            },
        }

        body_html = md.markdown(
            text,
            extensions=extensions,
            extension_configs=extension_configs,
        )

        pygments_css = HtmlFormatter(style="monokai").get_style_defs(".codehilite")

        document = f"""<!DOCTYPE html>
<html lang="en">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<style>
{_BASE_STYLES}
body {{ max-width: 860px; margin: 0 auto; }}
h1, h2, h3, h4, h5, h6 {{ color: #e0e0e0; margin-top: 1.4em; margin-bottom: 0.6em; }}
a {{ color: #569cd6; }}
code {{ background: #2d2d2d; padding: 2px 5px; border-radius: 3px; font-size: 0.9em; }}
pre {{ background: #2d2d2d; padding: 12px; border-radius: 6px; overflow-x: auto; }}
pre code {{ background: none; padding: 0; }}
blockquote {{ border-left: 3px solid #444; padding-left: 12px; color: #999; margin: 1em 0; }}
table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
th, td {{ border: 1px solid #444; padding: 8px 10px; text-align: left; }}
th {{ background: #2d2d2d; }}
img {{ max-width: 100%; }}
hr {{ border: none; border-top: 1px solid #444; margin: 1.5em 0; }}
{pygments_css}
</style></head>
<body>{body_html}</body></html>"""

        return document

    # ------------------------------------------------------------------
    # Preview type detection
    # ------------------------------------------------------------------

    @staticmethod
    def get_preview_type(content_type: str, filename: str) -> str | None:
        """Determine the preview type for a given content type and filename."""
        ct = content_type.lower() if content_type else ""
        ext = os.path.splitext(filename)[1].lower() if filename else ""

        # Office documents
        if ct in OFFICE_MIME_TYPES or ext in OFFICE_EXTENSIONS:
            if _has_libreoffice():
                return "pdf"
            if ext in _PYTHON_CONVERTIBLE:
                return "html"
            # Unsupported legacy format without LibreOffice
            return None

        if ct in CSV_MIME_TYPES or ext == ".csv":
            return "html"
        if ct in MARKDOWN_MIME_TYPES or ext in (".md", ".markdown"):
            return "html"
        if ct == "application/pdf":
            return "pdf"
        if ct.startswith("image/"):
            return "image"
        if ct.startswith("video/"):
            return "video"
        if ct.startswith("audio/"):
            return "audio"
        if ct.startswith("text/"):
            return "text"

        return None

    @staticmethod
    def can_convert_office(extension: str) -> bool:
        """Check if an office format can be converted in the current environment."""
        if _has_libreoffice():
            return extension.lower() in OFFICE_EXTENSIONS
        return extension.lower() in _PYTHON_CONVERTIBLE
